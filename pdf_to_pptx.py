import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import pdfplumber

# Set page title and description
st.set_page_config(page_title="PDF to PPTX Converter", layout="centered")

# Function to convert PDF to PPTX using pdfplumber for text extraction
def convert_pdf_to_pptx_(pdf_file):
    prs = Presentation()

    with pdfplumber.open(pdf_file) as pdf:
        for page_num in range(len(pdf.pages)):
            slide_layout = prs.slide_layouts[5]  # Use slide layout for content only
            slide = prs.slides.add_slide(slide_layout)
            content = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
            content.text = pdf.pages[page_num].extract_text()

    pptx_file = BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)
    return pptx_file

# Main function for the Streamlit app
def main():
    st.title("PDF to PPTX Converter")

    uploaded_file = st.file_uploader("Upload a PDF file", type=["pdf"])

    if uploaded_file is not None:
        st.write("PDF file uploaded successfully!")
        st.write("Converting PDF to PPTX...")

        pptx_file = convert_pdf_to_pptx(uploaded_file)

        # Create a download link for the converted PPTX file
        st.download_button(
            label="Download PPTX File",
            data=pptx_file,
            file_name=uploaded_file.name.split(".")[0] + ".pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

if __name__ == "__main__":
    main()
