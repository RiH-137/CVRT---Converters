import importlib
import subprocess
import streamlit as st
from pdf2docx import Converter
from io import BytesIO
import tempfile
import os
import sys


import pdfplumber


from docx import Document
from pptx import Presentation
from pptx.util import Inches



from PIL import Image
from PyPDF2 import PdfMerger

import time

# Set page title and description
st.set_page_config(page_title="CVRT - Converters", layout="centered")

# Home page
def home_page():
    st.title("CVRT - Converters")
    st.markdown("Welcome to CVRT! Choose an option from the sidebar to convert different file formats.")
  
    st.sidebar.title("Navigation")
    selection = st.sidebar.radio("Go to:", ["Home", "DOCX to PDF", "DOCX to PPTX", "Image/s to PDF", "JPEG to PNG", "PNG to JPEG", "PDF Merger", "PDF to DOCX", "PDF to PPTX", "About"])

    return selection
    
# About page
def about_page():
    st.title("About CVRT")
    st.markdown("""
    CVRT (Converters) is a Streamlit app that provides various file conversion functionalities. 
    You can convert DOCX to PDF, DOCX to PPTX, images to PDF, JPEG to PNG, PNG to JPEG, merge PDF files, 
    and convert PDF files to DOCX or PPTX formats. Choose an option from the sidebar to get started!
    """)

    st.title("About Me...")
    st.markdown("""
    Date:- 18 April 2024
    """)
    st.markdown("""
    🚀 Hey there! I'm Rishi, a passionate Computer Science & Engineering Undergraduate with a keen interest in the vast world of technology. Currently specializing in AI and Machine Learning, I'm on a perpetual quest for knowledge and thrive on learning new skills.

💻 My journey in the tech realm revolves around programming, problem-solving, and staying on the cutting edge of emerging technologies. With a strong foundation in Computer Science, I'm driven by the exciting intersection of innovation and research.

🔍 Amidst the digital landscape, I find myself delving into the realms of Blockchain, crafting Android Applications, and ML projects.
 JAVA and Python . 
My GitHub profile (https://github.com/RiH-137) showcases my ongoing commitment to refining my craft and contributing to the tech community.

🏎️ Outside the digital realm, I'm a fervent Formula 1 enthusiast, experiencing the thrill of high-speed pursuits. When I'm not immersed in code or cheering for my favorite F1 team, you might find me strategizing moves on the chessboard.

📧 Feel free to reach out if you're as passionate about technology as I am or if you just fancy a chat! You can connect with me at 101rishidsr@gmail.com and https://www.biodrop.io/RiH-137.

Let's build, innovate, and explore the limitless possibilities of technology together! 🌐✨
                
    """)
    st.image('1.png', caption="", use_column_width=True)
#--------------------------------------------------------------------------------------------------------------
# DOCX to PDF page
def docx_to_pdf_page():
    st.title("DOCX to PDF Converter")
    st.markdown("This page allows you to convert DOCX files to PDF format.")

    def convert_docx_to_pdf(docx_file):
    # save DOCX file to a temporary file
        temp_docx = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
        temp_docx.write(docx_file.read())
        temp_docx.close()

        # get the path for the temporary DOCX and PDF files
        temp_docx_path = temp_docx.name
        pdf_path = os.path.splitext(temp_docx_path)[0] + ".pdf"

        # use subprocess to call the docx2pdf CLI tool
        subprocess.run(["docx2pdf", temp_docx_path])

        return pdf_path


    # upload DOCX file with clear guidance
    docx_file = st.file_uploader("Upload a DOCX file to convert to PDF", type=['docx'])

    if docx_file is not None:
        # perform conversion when DOCX file is uploaded
        pdf_path = convert_docx_to_pdf(docx_file)

        # provide download link for the converted PDF file
        st.markdown("### Download Converted File")
        st.download_button(label="Download PDF", data=open(pdf_path, 'rb'), file_name="converted_document.pdf", mime="application/pdf")

    else:
        st.info("Please upload a DOCX file for conversion.")
    st.image('2.png', caption="", use_column_width=True)


#-------------------------------------------------------------------------------------------------


# DOCX to PPTX page
def docx_to_pptx_page():
    st.title("DOCX to PPTX Converter")
    st.markdown("This page allows you to convert DOCX files to PPTX format.")

    def convert_docx_to_pptx(docx_file):
        prs = Presentation()

        doc = Document(docx_file)
        
        for paragraph in doc.paragraphs:
            slide_layout = prs.slide_layouts[5]  # Use slide layout for content only
            slide = prs.slides.add_slide(slide_layout)
            content = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
            content.text = paragraph.text

        pptx_file = BytesIO()
        prs.save(pptx_file)
        pptx_file.seek(0)
        return pptx_file
    

    uploaded_docx_file = st.file_uploader("Upload a DOCX file", type=["docx"])

    if uploaded_docx_file is not None:
        st.write("DOCX file uploaded successfully!")
        st.write("Converting DOCX to PPTX...")

        pptx_file = convert_docx_to_pptx(uploaded_docx_file)

        # Create a download link for the converted PPTX file
        st.download_button(
            label="Download PPTX File",
            data=pptx_file,
            file_name=uploaded_docx_file.name.split(".")[0] + ".pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    st.image('3.png', caption="", use_column_width=True)









#----------------------------------------------------------------------------------------------------------
# Image to PDF page
def img_to_pdf_page():
    
    st.title("Image/s to PDF Converter")
    st.markdown("This page allows you to convert images to PDF format.")

    def convert_images_to_pdf(images):
        merger = PdfMerger()
        temp_files = []

        for img in images:
            img_bytes = img.read()  # Read the BytesIO object
            image = Image.open(BytesIO(img_bytes))

            pdf_bytes = BytesIO()
            image.save(pdf_bytes, format="PDF")
            pdf_bytes.seek(0)

            temp_files.append(pdf_bytes)

            merger.append(pdf_bytes)

        output_pdf = BytesIO()
        merger.write(output_pdf)

        for temp_file in temp_files:
            temp_file.close()

        return output_pdf
    
    uploaded_files = st.file_uploader("Upload images to generate PDF. You can select multiple images.", type=["jpg", "jpeg", "png"], accept_multiple_files=True)

    if uploaded_files:
        st.write("Images uploaded successfully!")
        st.write("Generating PDF...")

        pdf_file = convert_images_to_pdf(uploaded_files)

        #   generated PDF file
        st.download_button(
            label="Download PDF File",
            data=pdf_file.getvalue(),
            file_name="generated_pdf.pdf",
            mime="application/pdf",
        )
    st.image('4.png', caption="", use_column_width=True)        












#---------------------------------------------------------------------------------------------------------------
# JPEG to PNG page
def jpeg_to_png_page():
    st.title("JPEG to PNG Converter")
    st.markdown("This page allows you to convert JPEG images to PNG format.")

    def convert_jpeg_to_png(jpeg_file):
        image = Image.open(jpeg_file)
        png_file = BytesIO()
        image.save(png_file, format="PNG")
        png_file.seek(0)
        return png_file
    

    uploaded_file = st.file_uploader("Upload a JPEG file", type=["jpg", "jpeg"])

    if uploaded_file is not None:
        st.write("JPEG file uploaded successfully!")
        st.write("Converting JPEG to PNG...")

        png_file = convert_jpeg_to_png(uploaded_file)

        # Ddsplay the converted PNG images
        st.image(png_file, caption="Converted PNG Image")

        # create a download link for the converted PNG file
        st.download_button(
            label="Download PNG File",
            data=png_file,
            file_name=uploaded_file.name.split(".")[0] + ".png",
            mime="image/png",
        )
    st.image('5.png', caption="", use_column_width=True)    









#--------------------------------------------------------------------------------------------------------------------

# PNG to JPEG page
def png_to_jpeg_page():
    st.title("PNG to JPEG Converter")
    st.markdown("This page allows you to convert PNG images to JPEG format.")

    def convert_png_to_jpeg(png_file):
        image = Image.open(png_file)
        jpeg_file = BytesIO()
        image.convert("RGB").save(jpeg_file, format="JPEG")
        jpeg_file.seek(0)
        return jpeg_file

    uploaded_file = st.file_uploader("Upload a PNG file", type=["png"])

    if uploaded_file is not None:
            st.write("PNG file uploaded successfully!")
            st.write("Converting PNG to JPEG...")

            jpeg_file = convert_png_to_jpeg(uploaded_file)

            # sisplay the converted JPEG image
            st.image(jpeg_file, caption="Converted JPEG Image")

            # create a download link for the converted JPEG file
            st.download_button(
                label="Download JPEG File",
                data=jpeg_file,
                file_name=uploaded_file.name.split(".")[0] + ".jpg",
                mime="image/jpeg",
            )

    st.image('6.png', caption="", use_column_width=True)

        













#------------------------------------------------------------------------------------------------------------------------
# PDF Merger page
def pdf_merger_page():
    st.title("PDF Merger")
    st.markdown("This page allows you to merge multiple PDF files into one.")

    def merge_pdfs(pdfs):
        merger = PdfMerger()

        for pdf in pdfs:
            merger.append(pdf)

        merged_pdf = BytesIO()
        merger.write(merged_pdf)
        merged_pdf.seek(0)
        return merged_pdf

    uploaded_files = st.file_uploader("Upload PDF files to merge. (Note:- You can select multiple PDFs for merging.)", type=["pdf"], accept_multiple_files=True)

    if uploaded_files is not None:
        st.write("PDF files uploaded successfully!")
        st.write("Merging PDF files...")

        merged_pdf = merge_pdfs(uploaded_files)

        # create a download link for the merged PDF file
        st.download_button(
            label="Download Merged PDF",
            data=merged_pdf,
            file_name="merged_pdf.pdf",
            mime="application/pdf",
        )
    st.image('3.png', caption="", use_column_width=True)







#------------------------------------------------------------------------------------------------------------------
# PDF to DOCX page
def pdf_to_docx_page():
    st.title("PDF to DOCX Converter")
    st.markdown("This page allows you to convert PDF files to DOCX format.")

    def convert_pdf_to_docx(pdf_file):
        try:
            temp_pdf = tempfile.NamedTemporaryFile(delete=False)
            temp_pdf.write(pdf_file.read())
            temp_pdf.close()

            output_docx = BytesIO()

            cv = Converter(temp_pdf.name)
            cv.convert(output_docx, start=0, end=None)
            cv.close()

            os.unlink(temp_pdf.name)

            return output_docx

        except Exception as e:
            st.error(f"An error occurred during conversion: {str(e)}")
            st.error("Please check if the PDF file is valid and try again.")
            return None

    pdf_file = st.file_uploader("Upload a PDF file to convert to DOCX", type=['pdf'])

    if pdf_file is not None:
        output_docx = convert_pdf_to_docx(pdf_file)

        if output_docx is not None:
            st.markdown("### Download Converted File")
            st.download_button(label="Download DOCX", data=output_docx.getvalue(), file_name="converted_document.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    else:
        st.info("Please upload a PDF file for conversion.")





    st.image('2.png', caption="", use_column_width=True)











#-------------------------------------------------------------------------------------------------------------------    
# PDF to PPTX page


# PDF to PPTX page
def pdf_to_pptx_page():
    st.title("PDF to PPTX Converter")
    st.markdown("This page allows you to convert PDF files to PPTX format.")

    def convert_pdf_to_pptx(pdf_file):
        prs = Presentation()

        with pdfplumber.open(pdf_file) as pdf:
            for page_num in range(len(pdf.pages)):
                slide_layout = prs.slide_layouts[5]  # Use slide layout for content only
                slide = prs.slides.add_slide(slide_layout)
                content = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
                content.text = pdf.pages[page_num].extract_text()

        pptx_file = BytesIO()
        prs.save(pptx_file)
        pptx_file.seek(0)
        return pptx_file
    
    uploaded_file = st.file_uploader("Upload a PDF file", type=["pdf"])

    if uploaded_file is not None:
        st.write("PDF file uploaded successfully!")
        st.write("Converting PDF to PPTX...")

        pptx_file = convert_pdf_to_pptx(uploaded_file)

        # converted PPTX file
        st.download_button(
            label="Download PPTX File",
            data=pptx_file,
            file_name=uploaded_file.name.split(".")[0] + ".pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )


    st.image('5.png', caption="", use_column_width=True)


#-------------------------------------------------------------------------------------------------------------
#  streamlit application       !wooh!
def main():
    selection = home_page()

    
    if selection == "About":
        about_page()
    elif selection == "DOCX to PDF":
        docx_to_pdf_page()
    elif selection == "DOCX to PPTX":
        docx_to_pptx_page()
    elif selection == "Image/s to PDF":
        img_to_pdf_page()
    elif selection == "JPEG to PNG":
        jpeg_to_png_page()
    elif selection == "PNG to JPEG":
        png_to_jpeg_page()
    elif selection == "PDF Merger":
        pdf_merger_page()
    elif selection == "PDF to DOCX":
        pdf_to_docx_page()
    elif selection == "PDF to PPTX":
        pdf_to_pptx_page()

if __name__ == "__main__":
    main()
