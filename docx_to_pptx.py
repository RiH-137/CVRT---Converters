import streamlit as st
from io import BytesIO
from docx import Document
from pptx import Presentation
from pptx.util import Inches

# Set page title and description
st.set_page_config(page_title="DOCX to PPTX Converter", layout="centered")

# Function to convert DOCX to PPTX
def convert_docx_to_pptx(docx_file):
    prs = Presentation()

    doc = Document(docx_file)
    for paragraph in doc.paragraphs:
        slide_layout = prs.slide_layouts[5]  # Use slide layout for content only
        slide = prs.slides.add_slide(slide_layout)
        content = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
        content.text = paragraph.text

    pptx_file = BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)
    return pptx_file

# Main function for the Streamlit app
def main():
    st.title("DOCX to PPTX Converter")

    uploaded_docx_file = st.file_uploader("Upload a DOCX file", type=["docx"])

    if uploaded_docx_file is not None:
        st.write("DOCX file uploaded successfully!")
        st.write("Converting DOCX to PPTX...")

        pptx_file = convert_docx_to_pptx(uploaded_docx_file)

        # Create a download link for the converted PPTX file
        st.download_button(
            label="Download PPTX File",
            data=pptx_file,
            file_name=uploaded_docx_file.name.split(".")[0] + ".pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

if __name__ == "__main__":
    main()
